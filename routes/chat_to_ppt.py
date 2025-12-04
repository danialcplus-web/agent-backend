from fastapi import APIRouter, HTTPException
from pydantic import BaseModel
from typing import List, Optional
from services.supabase_client import supabase
from pptx import Presentation
from pptx.util import Pt
from starlette.responses import StreamingResponse
import io
import logging
from datetime import datetime
import json

logger = logging.getLogger(__name__)

router = APIRouter()


class MessageItem(BaseModel):
    role: str
    content: str


class ExportRequest(BaseModel):
    chat_id: Optional[str] = None
    messages: Optional[List[MessageItem]] = None

    class Config:
        extra = "ignore"



# -----------------------------------
# PPT CREATOR
# -----------------------------------
def create_ppt_from_chat(messages: List[dict], filename: Optional[str] = None) -> io.BytesIO:
    prs = Presentation()

    for msg in messages:
        try:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
        except Exception:
            slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Title
        try:
            title_shape = slide.shapes.title
            title_text = "User Message" if msg.get("role") == "user" else "Assistant Message"
            title_shape.text = title_text
        except Exception:
            pass

        # Body text
        body = None
        try:
            body = slide.placeholders[1].text_frame
        except Exception:
            for shp in slide.shapes:
                if hasattr(shp, "has_text_frame") and shp.has_text_frame:
                    body = shp.text_frame
                    break

        if body is None:
            from pptx.util import Inches
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
            body = txBox.text_frame

        # Clear any existing content
        body.clear()

        content = msg.get("content") or ""
        paragraphs = content.splitlines() if content else [""]

        for i, p in enumerate(paragraphs):
            if i == 0:
                para = body.paragraphs[0]
                para.text = p
            else:
                para = body.add_paragraph()
                para.text = p
            for run in para.runs:
                run.font.size = Pt(18)

    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio


# -----------------------------------
# API ROUTE
# -----------------------------------
@router.post("/chat/export/pptx")
async def export_chat_to_pptx(payload: ExportRequest):

    messages = []

    # Case 1 — direct message array
    if payload.messages:
        messages = [m.dict() for m in payload.messages]

    # Case 2 — load messages from Supabase using chat_id
    elif payload.chat_id:

        try:
            result = (
                supabase
                .from_("chats")
                .select("id, messages")
                .eq("id", payload.chat_id)
                .single()
                .execute()
            )
        except Exception as e:
            logger.exception("Supabase query failed: %s", e)
            raise HTTPException(status_code=500, detail="Failed to query chat row")

        # Access row data correctly
        data = result.data
        if not data:
            raise HTTPException(status_code=404, detail="Chat not found")

        msgs = data.get("messages")

        # If messages saved as a JSON string
        if isinstance(msgs, str):
            try:
                msgs = json.loads(msgs)
            except:
                msgs = []

        if not msgs:
            raise HTTPException(status_code=404, detail="No messages found for chat_id")

        # Normalize to array of {role, content}
        for m in msgs:
            if isinstance(m, dict) and ("role" in m and "content" in m):
                messages.append({
                    "role": m["role"],
                    "content": m.get("content", "")
                })


    else:
        raise HTTPException(status_code=400, detail="Provide either chat_id or messages")

    # Build output PPT
    filename = f"chat_export_{datetime.utcnow().strftime('%Y%m%dT%H%M%SZ')}.pptx"
    ppt_bytes = create_ppt_from_chat(messages, filename=filename)

    return StreamingResponse(
        ppt_bytes,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={
            "Content-Disposition": f'attachment; filename="{filename}"'
        }
    )
