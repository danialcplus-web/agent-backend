# back_end/services/file_processing.py
import io
import logging
from typing import Optional, Tuple, List
import fitz  # PyMuPDF
import docx
import pandas as pd
import pptx

logger = logging.getLogger(__name__)

def extract_text_from_pdf_bytes(b: bytes) -> str:
    text_parts = []
    try:
        doc = fitz.open(stream=b, filetype="pdf")
        for page in doc:
            text = page.get_text("text")
            if text:
                text_parts.append(text)
    except Exception as e:
        logger.exception("PDF extraction failed: %s", e)
    return "\n\n".join(text_parts)

def extract_text_from_docx_bytes(b: bytes) -> str:
    try:
        bio = io.BytesIO(b)
        doc = docx.Document(bio)
        paragraphs = [p.text for p in doc.paragraphs if p.text]
        return "\n\n".join(paragraphs)
    except Exception as e:
        logger.exception("DOCX extraction failed: %s", e)
        return ""

def extract_text_from_txt_bytes(b: bytes) -> str:
    try:
        return b.decode("utf-8", errors="ignore")
    except Exception:
        return str(b)

def extract_text_from_csv_bytes(b: bytes) -> str:
    try:
        df = pd.read_csv(io.BytesIO(b))
        rows = []
        for i, r in df.iterrows():
            rows.append(" | ".join([f"{c}:{r[c]}" for c in df.columns]))
        return "\n".join(rows)
    except Exception:
        try:
            return b.decode("utf-8", errors="ignore")
        except:
            return ""
        
def extract_text_from_pptx_bytes(b: bytes) -> str:
    try:
        bio = io.BytesIO(b)
        presentation = pptx.Presentation(bio)
        text_runs = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n\n".join(text_runs)
    except Exception as e:
        logger.exception("PPTX extraction failed: %s", e)
        return ""

def extract_text_from_file_bytes(filename: str, content_bytes: Optional[bytes] = None ) -> str:
    #fname = filename.filename
    lower = filename.lower()
    if lower.endswith(".pdf"):
        return extract_text_from_pdf_bytes(content_bytes)
    if lower.endswith(".docx") or lower.endswith(".doc"):
        return extract_text_from_docx_bytes(content_bytes)
    if lower.endswith(".csv"):
        return extract_text_from_csv_bytes(content_bytes)
    if lower.endswith(".txt"):
        return extract_text_from_txt_bytes(content_bytes)
    if lower.endswith(".pptx") or lower.endswith(".ppt"):
        return extract_text_from_pptx_bytes(content_bytes)
    # default
    try:
        return content_bytes.decode("utf-8", errors="ignore")
    except:
        return ""